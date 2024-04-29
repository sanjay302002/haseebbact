import streamlit as st
from streamlit_option_menu import option_menu
#from PyPDF2 import PdfReader
#from pptx import Presentation
#from docx import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
import os
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain_community.vectorstores import FAISS
from dotenv import load_dotenv

st.set_page_config(layout="wide")

load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

def get_pdf_text(pdf_docs):
    from PyPDF2 import PdfReader
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_ppt_text(ppt_files):
    from pptx import Presentation
    text = ""
    for ppt_file in ppt_files:
        presentation = Presentation(ppt_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    return text

def get_docx_text(docx_files):
    from docx import Document
    text = ""
    try:
        for docx_file in docx_files:
            doc = Document(docx_file)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
    except Exception as e:
        st.error(f"Error occurred while processing Word document: {e}")
        return None
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    st.write("Reply: ", response["output_text"])

#Horizontal
#to make it sidebar use (with st.sidebar:) & also erase (orientation="horizontal")
with st.sidebar:
     selected = option_menu('Multiple Document Querying System',
                                    ['Ask with PDF',
                                     'Ask with PPT',
                                     'Ask with Docx'],
                                     default_index=0,
                                     
                          )

def main():
    if selected == 'Ask with PDF':
        st.title('Chat with PDF using Gemini')
        
        pdf_docs = st.file_uploader("Upload your Portable Document Files and Click on the Submit & Process Button", accept_multiple_files=True)
        if st.button("Submit & Process"):
                with st.spinner("Processing..."):
                    raw_text = get_pdf_text(pdf_docs)
                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks)
                    st.success("Done")
        user_question = st.text_input("Ask a Question from the PDF Files")
        if user_question:
            user_input(user_question)
       
    elif selected == 'Ask with PPT':
        st.title('Chat with PPT using Gemini')    
        ppt_files = st.file_uploader("Upload your PowerPoint Files and Click on the Submit & Process Button", accept_multiple_files=True)
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                raw_text = get_ppt_text(ppt_files)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("Done")
        user_question = st.text_input("Ask a Question from the PPT Files")
        if user_question:
            user_input(user_question)
            # Your code for PPT interaction here
    elif selected == 'Ask with Docx':
        st.title('Chat with DOCX using Gemini')
        docx_files = st.file_uploader("Upload your Word Document and Click on the Submit & Process Button", accept_multiple_files=True)
       
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                raw_text = get_docx_text(docx_files)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("Done")
        user_question = st.text_input("Ask a Question from the Word Document")
        if user_question:
            user_input(user_question)
           
        # Your code for DOCX interaction here

if __name__ == "__main__":
    main()
